#!/usr/bin/env python3
"""
Minimal black & white deck for the WPU talk, with a dedicated sources slide
(citations match web/data/factcheck.ts + app methodology).

Run:  .venv-pptx/bin/python presentation/build_wpu_deck.py
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "AI_Fact_Checker_WPU_10min.pptx"
APP_URL = "https://web-tau-peach-63.vercel.app"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)  # body secondary
MUTED = RGBColor(0x6E, 0x6E, 0x6E)  # URLs / metadata


@dataclass(frozen=True)
class SourceLine:
    title: str
    url: str
    kind: str  # e.g. Preprint, Institution, Company


# Same ground-truth list as the web app’s AI-detection fact checks + methodology.
SOURCES: list[SourceLine] = [
    SourceLine(
        "GPT detectors are biased against non-native English writers (Liang et al., 2023)",
        "https://arxiv.org/abs/2304.02819",
        "Preprint",
    ),
    SourceLine(
        "Can AI-Generated Text be Reliably Detected? (Sadasivan et al., 2023)",
        "https://arxiv.org/abs/2303.11156",
        "Preprint",
    ),
    SourceLine(
        "New AI Classifier for Indicating AI-Written Text — OpenAI (discontinued, 2023)",
        "https://openai.com/index/new-ai-classifier-for-indicating-ai-written-text/",
        "Company",
    ),
    SourceLine(
        "Guidance on AI detection; disabling Turnitin’s AI detector (Vanderbilt, 2023)",
        "https://news.vanderbilt.edu/2023/08/16/guidance-on-ai-detection-and-why-were-disabling-turnitins-ai-detector/",
        "Institution",
    ),
    SourceLine(
        "AI writing detection — update from Turnitin’s CPO (Turnitin, 2023)",
        "https://www.turnitin.com/blog/ai-writing-detection-update-from-turnitins-chief-product-officer",
        "Company",
    ),
    SourceLine(
        "All app content: model cards, benchmarks, and fact-checks cite dated primary sources (docs, papers, or vendor materials).",
        APP_URL,
        "App",
    ),
]


def _bg_white(slide) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = WHITE


def _h_rule(slide, y_in: float) -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y_in), Inches(8.6), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()


def _add_title_block(slide, y: float, text: str, size: int = 28) -> None:
    t = slide.shapes.add_textbox(Inches(0.85), Inches(y), Inches(8.3), Inches(0.75))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT


def _body(slide, y: float, h: float, lines: list[str], size: int = 15) -> None:
    t = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(8.2), Inches(h))
    tf = t.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY if not line.strip().startswith("•") else GRAY
        p.space_after = Pt(8)
        p.line_spacing = 1.2


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def _add_source_column(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[SourceLine],
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, s in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(10)
        r0 = p.add_run()
        r0.text = s.title + "\n"
        r0.font.size = Pt(11.5)
        r0.font.bold = True
        r0.font.color.rgb = BLACK
        r1 = p.add_run()
        r1.text = s.url + "\n"
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = MUTED
        r2 = p.add_run()
        r2.text = s.kind
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = MUTED
        r2.font.italic = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # —— 1. Title ——
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_white(s0)
    _h_rule(s0, 0.45)
    t0 = s0.shapes.add_textbox(Inches(0.9), Inches(1.15), Inches(8.2), Inches(1.2))
    t0f = t0.text_frame
    t0f.text = "AI Fact Checker"
    t0f.paragraphs[0].font.size = Pt(40)
    t0f.paragraphs[0].font.bold = True
    t0f.paragraphs[0].font.color.rgb = BLACK
    t0f.paragraphs[0].alignment = PP_ALIGN.LEFT
    s1 = t0f.add_paragraph()
    s1.text = "Cade Kukk  ·  West Pennsylvania University"
    s1.font.size = Pt(16)
    s1.font.color.rgb = GRAY
    s1.space_before = Pt(14)
    s2 = t0f.add_paragraph()
    s2.text = "~10 minutes  ·  verified, source-first reference on AI"
    s2.font.size = Pt(13)
    s2.font.color.rgb = MUTED
    s2.space_before = Pt(10)
    # Footer line
    _h_rule(s0, 4.9)
    foot = s0.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(8.2), Inches(0.4))
    ff = foot.text_frame
    ff.text = "black & white deck  ·  sources on next slides"
    ff.paragraphs[0].font.size = Pt(9)
    ff.paragraphs[0].font.color.rgb = MUTED
    _notes(s0, "Introduce yourself, thank the host, state you will end with a clear source list for the room.")

    # —— 2. What it is (minimal) ——
    s1_ = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_white(s1_)
    _add_title_block(s1_, 0.5, "What it is", 24)
    _h_rule(s1_, 0.95)
    _body(
        s1_,
        1.2,
        2.0,
        [
            "A free web reference: major models, clear definitions, and fact checks — every claim tied to a primary source.",
            "It is not an AI text detector. It is the opposite: look up what models are, with citations, before judging text.",
        ],
        16,
    )
    _notes(s1_, "Two sentences max. Then pivot: the story and the sources matter because detectors are shaky (listed next).")

    # —— 3. Sources (two columns, visually balanced) ——
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_white(s2)
    # Vertical accent
    bar = s2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.06), Inches(4.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

    _add_title_block(s2, 0.45, "Sources (selected)", 24)
    sub = s2.shapes.add_textbox(Inches(0.85), Inches(0.9), Inches(8.3), Inches(0.4))
    sf = sub.text_frame
    sf.text = "The same materials ground the in-app fact checks; use these for your own follow-up or syllabus."
    sf.paragraphs[0].font.size = Pt(11.5)
    sf.paragraphs[0].font.color.rgb = GRAY

    mid = (len(SOURCES) + 1) // 2
    col_a = SOURCES[:mid]
    col_b = SOURCES[mid:]
    _add_source_column(s2, 0.75, 1.35, 4.0, 3.7, col_a)
    _add_source_column(s2, 5.15, 1.35, 4.0, 3.7, col_b)

    _h_rule(s2, 4.75)
    cap = s2.shapes.add_textbox(Inches(0.85), Inches(4.8), Inches(8.3), Inches(0.5))
    cf = cap.text_frame
    cf.text = "arXiv.org  ·  OpenAI  ·  Vanderbilt  ·  Turnitin  —  + live app: " + APP_URL
    cf.paragraphs[0].font.size = Pt(9.5)
    cf.paragraphs[0].font.color.rgb = MUTED
    _notes(s2, "Point to columns. If asked “where is this from?”, you already showed it. Offer to share the deck or URL.")

    # —— 4. QR + thank you ——
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_white(s3)
    _add_title_block(s3, 0.5, "Try the app", 24)
    _h_rule(s3, 0.95)

    qr = qrcode.make(APP_URL, box_size=4, border=1)
    buf = io.BytesIO()
    qr.save(buf, format="PNG")
    buf.seek(0)
    # Center QR
    s3.shapes.add_picture(buf, Inches(3.55), Inches(1.3), width=Inches(2.7))

    u = s3.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(7.0), Inches(0.4))
    uf = u.text_frame
    uf.text = APP_URL
    uf.paragraphs[0].font.size = Pt(13)
    uf.paragraphs[0].font.color.rgb = BLACK
    uf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tq = s3.shapes.add_textbox(Inches(1.2), Inches(4.45), Inches(7.6), Inches(0.55))
    tqf = tq.text_frame
    tqf.text = "Free  ·  no login  ·  questions welcome"
    tqf.paragraphs[0].font.size = Pt(12)
    tqf.paragraphs[0].font.color.rgb = GRAY
    tqf.paragraphs[0].alignment = PP_ALIGN.CENTER
    th = s3.shapes.add_textbox(Inches(1.2), Inches(4.95), Inches(7.6), Inches(0.4))
    thf = th.text_frame
    thf.text = "Thank you"
    thf.paragraphs[0].font.size = Pt(18)
    thf.paragraphs[0].font.bold = True
    thf.paragraphs[0].font.color.rgb = BLACK
    thf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _notes(s3, "Hold for applause. If tech fails, read URL aloud. Q&A: data is curated; updated as the field changes.")

    prs.save(OUT)
    print(f"Wrote: {OUT} (4 slides, minimal B&W, sources on slide 3)")


if __name__ == "__main__":
    main()
